"""
Document Processor - Multi-format document parsing and chunking
"""

import os
import re
import io
from typing import List, Dict, Any, Union
from dotenv import load_dotenv

load_dotenv()


class DocumentProcessor:
    """Process documents into chunks. Supports Markdown, PDF, Word, and PowerPoint."""

    def __init__(self, chunk_size: int = None, chunk_overlap: int = None):
        self.chunk_size = chunk_size or int(os.environ.get("CHUNK_SIZE", "512"))
        self.chunk_overlap = chunk_overlap or int(os.environ.get("CHUNK_OVERLAP", "50"))

    def _clean_text(self, text: str) -> str:
        """Clean and normalize text."""
        # Remove excessive whitespace
        text = re.sub(r"\n{3,}", "\n\n", text)
        text = re.sub(r" {2,}", " ", text)
        return text.strip()

    def _extract_pdf_text(self, file_content: bytes) -> str:
        """
        Extract text from PDF file.

        Args:
            file_content: PDF file content as bytes

        Returns:
            Extracted text content
        """
        try:
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(file_content))
            text_parts = []

            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)

            return "\n\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"Failed to extract PDF text: {str(e)}")

    def _extract_docx_text(self, file_content: bytes) -> str:
        """
        Extract text from Word document.

        Args:
            file_content: Word file content as bytes

        Returns:
            Extracted text content
        """
        try:
            from docx import Document

            doc = Document(io.BytesIO(file_content))
            text_parts = []

            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)

            return "\n\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"Failed to extract Word text: {str(e)}")

    def _extract_pptx_text(self, file_content: bytes) -> str:
        """
        Extract text from PowerPoint presentation.

        Args:
            file_content: PowerPoint file content as bytes

        Returns:
            Extracted text content
        """
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(file_content))
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())

                if slide_texts:
                    text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))

            return "\n\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"Failed to extract PPT text: {str(e)}")

    def extract_text_from_file(
        self, file_content: Union[bytes, str], filename: str
    ) -> str:
        """
        Extract text from various file formats.

        Args:
            file_content: File content (bytes for binary files, str for text files)
            filename: Original filename with extension

        Returns:
            Extracted text content

        Raises:
            ValueError: If file format is not supported
        """
        ext = os.path.splitext(filename.lower())[1]

        if ext == ".md" or ext == ".txt":
            # Text files - content should be string
            if isinstance(file_content, bytes):
                return file_content.decode("utf-8")
            return file_content

        elif ext == ".pdf":
            # PDF files - content should be bytes
            if isinstance(file_content, str):
                raise ValueError("PDF content must be bytes")
            return self._extract_pdf_text(file_content)

        elif ext == ".docx":
            # Word files - content should be bytes
            if isinstance(file_content, str):
                raise ValueError("Word content must be bytes")
            return self._extract_docx_text(file_content)

        elif ext == ".pptx":
            # PowerPoint files - content should be bytes
            if isinstance(file_content, str):
                raise ValueError("PPT content must be bytes")
            return self._extract_pptx_text(file_content)

        else:
            raise ValueError(f"Unsupported file format: {ext}")

    def _split_into_chunks(self, text: str, source_file: str) -> List[Dict[str, Any]]:
        """Split text into overlapping chunks."""
        text = self._clean_text(text)

        if len(text) <= self.chunk_size:
            return [{"text": text, "source_file": source_file, "chunk_index": 0}]

        chunks = []
        start = 0
        chunk_index = 0

        while start < len(text):
            end = start + self.chunk_size

            # Try to break at sentence boundary
            if end < len(text):
                # Look for sentence ending punctuation
                last_period = text.rfind("。", start, end)
                last_newline = text.rfind("\n", start, end)
                last_space = text.rfind(" ", start, end)

                # Prefer breaking at sentence or paragraph boundary
                break_point = max(last_period, last_newline)
                if break_point > start + self.chunk_size // 2:
                    end = break_point + 1
                elif last_space > start + self.chunk_size // 2:
                    end = last_space

            chunk_text = text[start:end].strip()

            if chunk_text:
                chunks.append(
                    {
                        "text": chunk_text,
                        "source_file": source_file,
                        "chunk_index": chunk_index,
                    }
                )
                chunk_index += 1

            # Move start with overlap
            start = end - self.chunk_overlap
            if start < 0:
                start = end

        return chunks

    def process_file(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Process a single Markdown file.

        Args:
            file_path: Path to the Markdown file

        Returns:
            List of chunk dictionaries with text and metadata
        """
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()

        filename = os.path.basename(file_path)
        return self._split_into_chunks(content, filename)

    def process_directory(self, dir_path: str) -> List[Dict[str, Any]]:
        """
        Process all Markdown files in a directory.

        Args:
            dir_path: Path to directory containing Markdown files

        Returns:
            List of all chunks from all files
        """
        all_chunks = []

        for filename in os.listdir(dir_path):
            if filename.endswith(".md"):
                file_path = os.path.join(dir_path, filename)
                chunks = self.process_file(file_path)
                all_chunks.extend(chunks)

        return all_chunks

    def process_uploaded_content(
        self, content: str, filename: str
    ) -> List[Dict[str, Any]]:
        """
        Process uploaded file content directly.

        Args:
            content: File content as string
            filename: Original filename

        Returns:
            List of chunk dictionaries
        """
        return self._split_into_chunks(content, filename)

    def process_uploaded_file(
        self, file_content: Union[bytes, str], filename: str
    ) -> List[Dict[str, Any]]:
        """
        Process uploaded file with automatic format detection.

        Supports: .md, .txt, .pdf, .docx, .pptx

        Args:
            file_content: File content (bytes for binary files, str for text files)
            filename: Original filename with extension

        Returns:
            List of chunk dictionaries
        """
        text = self.extract_text_from_file(file_content, filename)
        return self._split_into_chunks(text, filename)
